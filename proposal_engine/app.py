import os
import json
import zipfile
from io import BytesIO
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

TEMPLATE_PATH = r"C:\Users\GuillaumeTant\Releaf Carbon\Releaf Carbon - Documents\2. Commercial\6. Template\Template master proposition v3.pptx"
CONFIG_PATH = os.path.join(os.path.dirname(__file__), "slide_config.json")
PORT = 5001

# Mapping mission display name → clés dans slide_config.json
MISSION_MAP = {
    "Bilan Carbone": {
        "section":  "Bilan_Carbone",
        "cal":      "Bilan Carbone",
        "fin":      "Bilan_Carbone",
        "intitule": "Mesure de l'empreinte carbone",
    },
    "ACV": {
        "section":  "ACV",
        "cal":      "ACV",
        "fin":      "ACV",
        "intitule": "Analyse de Cycle de Vie",
    },
    "FDES / PEP": {
        "section":  "FDES_PEP",
        "cal":      "FDES_PEP",
        "fin":      "FDES_PEP",
        "intitule": "FDES / PEP",
    },
    "EPD": {
        "section":  "EPD",
        "cal":      "EPD",
        "fin":      "EPD",
        "intitule": "Environmental Product Declaration",
    },
}

# Mapping clé de subvention → label humain + variables à pré-remplir
SUBVENTION_CONFIG = {
    "Rev3_50pct": {
        "label":      "Booster Transformation – Rev3 (50%)",
        "programme":  "Booster Transformation",
        "operateur":  "Rev3",
        "pct":        "50%",
    },
    "BPI_40pct": {
        "label":      "Diag Décarbon'Action – Bpifrance (40%)",
        "programme":  "Diag Décarbon\u2019Action",
        "operateur":  "Bpifrance",
        "pct":        "40%",
    },
    "Rev3_30pct": {
        "label":      "Booster Transformation – Rev3 (30%)",
        "programme":  "Booster Transformation",
        "operateur":  "Rev3",
        "pct":        "30%",
    },
    "BPI_70pct": {
        "label":      "Diag Ecoconception – Bpifrance (70%)",
        "programme":  "Diag Ecoconception",
        "operateur":  "Bpifrance",
        "pct":        "70%",
    },
    "BPI_60pct": {
        "label":      "Diag Ecoconception – Bpifrance (60%)",
        "programme":  "Diag Ecoconception",
        "operateur":  "Bpifrance",
        "pct":        "60%",
    },
    "standard": {
        "label":      "Sans subvention",
        "programme":  "",
        "operateur":  "",
        "pct":        "",
    },
}

# Missions sans options de subvention (slide unique fixe)
MISSIONS_NO_SUBVENTION = {"FDES / PEP", "EPD"}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def load_config():
    if not os.path.exists(CONFIG_PATH):
        raise FileNotFoundError(f"slide_config.json introuvable : {CONFIG_PATH}")
    with open(CONFIG_PATH, encoding="utf-8") as f:
        return json.load(f)


def get_subventions_for_mission(mission, config):
    """Retourne la liste des clés de subvention disponibles pour une mission."""
    if mission in MISSIONS_NO_SUBVENTION:
        return ["standard"]
    m = MISSION_MAP[mission]
    fin_key = m["fin"]
    fin_section = config["sections"]["proposition_financiere"]["slides_per_mission"]
    entry = fin_section.get(fin_key)
    if entry is None:
        return ["standard"]
    if isinstance(entry, dict) and "options" in entry:
        return list(entry["options"].keys())
    return ["standard"]


def get_financial_slide(mission, subvention_key, config):
    """Retourne l'indice 1-based de la slide financière."""
    m = MISSION_MAP[mission]
    fin_key = m["fin"]
    fin_section = config["sections"]["proposition_financiere"]["slides_per_mission"]
    entry = fin_section.get(fin_key)
    if entry is None:
        raise ValueError(f"Aucune slide financière pour la mission '{mission}'")
    if isinstance(entry, int):
        return entry
    if isinstance(entry, dict) and "options" in entry:
        opts = entry["options"]
        if subvention_key not in opts:
            raise ValueError(f"Clé de subvention inconnue : {subvention_key}")
        return opts[subvention_key]
    raise ValueError(f"Format inattendu pour la mission '{mission}'")


def compute_slides_to_keep(mission, subvention_key, config):
    """Retourne un set d'indices 0-based des slides à conserver."""
    m = MISSION_MAP[mission]
    keep = set()

    # Introduction (toujours slides 1-9)
    for s in config["sections"]["introduction"]["slides"]:
        keep.add(s - 1)

    # Section mission
    section_key = m["section"]
    for s in config["sections"][section_key]["slides"]:
        keep.add(s - 1)

    # Calendrier header + slide mission
    cal_section = config["sections"]["calendrier"]
    keep.add(cal_section["section_header_slide"] - 1)
    cal_slide = cal_section["slides_per_mission"][m["cal"]]
    keep.add(cal_slide - 1)

    # Proposition financière header + slide sélectionnée
    pf_section = config["sections"]["proposition_financiere"]
    keep.add(pf_section["section_header_slide"] - 1)
    financial_slide = get_financial_slide(mission, subvention_key, config)
    keep.add(financial_slide - 1)

    return keep


def delete_slides(prs, keep_indices):
    """Supprime les slides dont l'indice n'est pas dans keep_indices."""
    sldIdLst = prs.slides._sldIdLst
    total = len(prs.slides)
    to_remove = sorted([i for i in range(total) if i not in keep_indices], reverse=True)
    for i in to_remove:
        sldIdLst.remove(sldIdLst[i])


def replace_text_zip(buf_in, replacements):
    """Remplace les placeholders {{VAR}} dans les XML du PPTX via ZIP."""
    buf_out = BytesIO()
    with zipfile.ZipFile(buf_in, "r") as zin, \
         zipfile.ZipFile(buf_out, "w", zipfile.ZIP_DEFLATED) as zout:
        for name in zin.namelist():
            data = zin.read(name)
            if name.endswith(".xml") or name.endswith(".rels"):
                text = data.decode("utf-8")
                for placeholder, value in replacements.items():
                    text = text.replace(placeholder, str(value))
                data = text.encode("utf-8")
            zout.writestr(name, data)
    buf_out.seek(0)
    return buf_out


def replace_logo(buf_in, logo_bytes_io):
    """Remplace la dernière image du slide 1 par le logo client."""
    prs = Presentation(buf_in)
    slide = prs.slides[0]
    pics = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    if pics:
        last_pic = pics[-1]
        left   = last_pic.left
        top    = last_pic.top
        width  = last_pic.width
        height = last_pic.height
        sp = last_pic._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(logo_bytes_io, left, top, width, height)
    buf_out = BytesIO()
    prs.save(buf_out)
    buf_out.seek(0)
    return buf_out


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/config")
def get_config():
    """Retourne la liste des missions et leurs subventions disponibles."""
    try:
        config = load_config()
        result = {}
        for mission in MISSION_MAP:
            subs = get_subventions_for_mission(mission, config)
            result[mission] = [
                {"key": k, "label": SUBVENTION_CONFIG.get(k, {}).get("label", k)}
                for k in subs
            ]
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/generate", methods=["POST"])
def generate():
    try:
        nom_entreprise  = (request.form.get("nom_entreprise") or "").strip()
        mission         = request.form.get("mission", "")
        subvention_key  = request.form.get("subvention", "standard")
        logo_file       = request.files.get("logo")

        if not nom_entreprise:
            return jsonify({"error": "Nom de l'entreprise requis"}), 400
        if mission not in MISSION_MAP:
            return jsonify({"error": f"Mission inconnue : {mission}"}), 400

        config = load_config()

        # 1. Calculer les slides à garder
        keep = compute_slides_to_keep(mission, subvention_key, config)

        # 2. Ouvrir le template, supprimer les slides inutiles
        if not os.path.exists(TEMPLATE_PATH):
            return jsonify({"error": f"Template PPTX introuvable : {TEMPLATE_PATH}"}), 500

        prs = Presentation(TEMPLATE_PATH)
        delete_slides(prs, keep)

        buf1 = BytesIO()
        prs.save(buf1)
        buf1.seek(0)

        # 3. Construire les variables de remplacement
        sub_info = SUBVENTION_CONFIG.get(subvention_key, SUBVENTION_CONFIG["standard"])
        m_info   = MISSION_MAP[mission]

        replacements = {
            "{{NOM_ENTREPRISE}}":        nom_entreprise,
            "{{TYPE_MISSION}}":          mission,
            "{{PROGRAMME_SUBVENTION}}":  sub_info["programme"],
            "{{OPERATEUR_SUBVENTION}}":  sub_info["operateur"],
            "{{POURCENTAGE_SUBVENTION}}": sub_info["pct"],
            "{{MONTANT_SUBVENTION}}":    "",   # baked in the slide
            "{{PRIX_APRES_SUBVENTION}}": "",   # baked in the slide
            "{{INTITULE_MISSION}}":      m_info["intitule"],
        }

        # 4. Remplacement texte via ZIP
        buf2 = replace_text_zip(buf1, replacements)

        # 5. Remplacement logo si fourni
        if logo_file and logo_file.filename:
            logo_bytes_io = BytesIO(logo_file.read())
            buf2 = replace_logo(buf2, logo_bytes_io)

        # 6. Retourner le fichier
        safe_nom = "".join(c for c in nom_entreprise if c.isalnum() or c in " _-").strip()
        filename = f"Proposition_{safe_nom}.pptx"
        buf2.seek(0)
        return send_file(
            buf2,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=filename,
        )

    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=PORT, debug=True)
